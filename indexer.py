import json
import os
import re
import sqlite3
import time
import unicodedata
import yaml
from multiprocessing import Process, Queue
import config

try:
    from langchain_text_splitters import RecursiveCharacterTextSplitter
except ImportError:
    from langchain.text_splitter import RecursiveCharacterTextSplitter

from langchain_chroma import Chroma
from langchain_core.documents import Document

from embedding_cache import CachedEmbeddings

import warnings
warnings.filterwarnings("ignore", category=UserWarning)

import logging
logging.getLogger("pypdf").setLevel(logging.ERROR)

import pypdf
import docx
import openpyxl
from pptx import Presentation

FILE_READ_TIMEOUT = 30
MAX_FILE_SIZE_MB = 100
XLSX_SKIP_SIZE_MB = 10

ZIP_MAGIC = b'PK\x03\x04'

# DashScope text-embedding-v2 单行最大 2048 Token，混合文本可能 1 字>1 token，保守截断
EMBED_MAX_CHARS = 600

_splitter = RecursiveCharacterTextSplitter(
    chunk_size=config.CHUNK_SIZE,
    chunk_overlap=config.CHUNK_OVERLAP,
)


def _ensure_dashscope_api_key():
    if os.environ.get("DASHSCOPE_API_KEY"):
        return True
    if getattr(config, "MY_API_KEY", ""):
        os.environ["DASHSCOPE_API_KEY"] = config.MY_API_KEY
        return True
    print("❌ 未配置 DashScope API Key。请设置环境变量 DASHSCOPE_API_KEY 或在 config.py 中填写 MY_API_KEY。")
    return False


def _truncate_for_embed(text: str) -> str:
    """确保文本长度不超过 Embedding API 限制，空串返回空格"""
    if not text or not text.strip():
        return " "
    if len(text) > EMBED_MAX_CHARS:
        return text[:EMBED_MAX_CHARS]
    return text


def normalize_path(path_str):
    """将 macOS 常见的 NFD 编码强制转为 NFC 标准编码"""
    return unicodedata.normalize('NFC', path_str)


def _init_scan_cache(conn: sqlite3.Connection):
    conn.execute("""
        CREATE TABLE IF NOT EXISTS scan_cache (
            filepath TEXT PRIMARY KEY,
            mtime REAL,
            source_type TEXT,
            chunks_json TEXT
        )
    """)
    conn.commit()


def _load_cached_docs(conn: sqlite3.Connection, filepath: str, mtime: float) -> list[Document] | None:
    """若缓存存在且 mtime 匹配，返回缓存的 Document 列表，否则返回 None"""
    row = conn.execute(
        "SELECT mtime, chunks_json FROM scan_cache WHERE filepath = ?",
        (filepath,),
    ).fetchone()
    if row is None:
        return None
    cached_mtime, chunks_json = row
    # mtime 是浮点，文件系统/SQLite 读写可能带来极小误差
    if cached_mtime is None or abs(float(cached_mtime) - float(mtime)) > 0.01:
        return None
    try:
        chunks = json.loads(chunks_json)
        return [Document(page_content=c["page_content"], metadata=c["metadata"]) for c in chunks]
    except (json.JSONDecodeError, KeyError):
        return None


def _save_cached_docs(conn: sqlite3.Connection, filepath: str, mtime: float, source_type: str, docs: list[Document], *, auto_commit: bool = True):
    chunks = [{"page_content": d.page_content, "metadata": d.metadata} for d in docs]
    conn.execute(
        "INSERT OR REPLACE INTO scan_cache (filepath, mtime, source_type, chunks_json) VALUES (?, ?, ?, ?)",
        (filepath, mtime, source_type, json.dumps(chunks, ensure_ascii=False)),
    )
    if auto_commit:
        conn.commit()


def _prune_scan_cache(conn: sqlite3.Connection, valid_filepaths: set[str]):
    """删除已不存在于磁盘的文件的缓存条目（批量 DELETE 提升性能）"""
    cursor = conn.execute("SELECT filepath FROM scan_cache")
    to_delete = [row[0] for row in cursor.fetchall() if row[0] not in valid_filepaths]
    if not to_delete:
        return
    batch_size = 500  # SQLite IN 子句限制
    for i in range(0, len(to_delete), batch_size):
        batch = to_delete[i : i + batch_size]
        placeholders = ",".join("?" * len(batch))
        conn.execute(f"DELETE FROM scan_cache WHERE filepath IN ({placeholders})", batch)
    conn.commit()


def build_documents_for_path_cached(
    filepath: str, mtime: float, source_type: str, conn: sqlite3.Connection | None
) -> list[Document]:
    """为单个文件构建 Document，优先使用扫描缓存。供全量/增量索引共用。"""
    if conn:
        cached = _load_cached_docs(conn, filepath, mtime)
        if cached is not None:
            return cached
    if source_type == "file":
        filename = os.path.basename(filepath)
        _, ext = os.path.splitext(filename)
        ext = ext.lower()
        docs = build_documents_for_file(filepath, filename, ext, source_type="file")
    else:
        try:
            with open(filepath, "r", encoding="utf-8", errors="ignore") as f:
                content = f.read()
        except OSError:
            return []
        docs = build_documents_for_mweb(filepath, content)
    if conn and docs:
        _save_cached_docs(conn, filepath, mtime, source_type, docs)
    return docs


def _read_file_worker(filepath: str, ext: str, q: Queue):
    """在子进程中执行文件读取，返回 (content, headings)，可被父进程强制终止"""
    warnings.filterwarnings("ignore", category=UserWarning)
    logging.getLogger("pypdf").setLevel(logging.ERROR)
    text = ""
    headings = []
    try:
        if ext in config.TEXT_EXTENSIONS:
            with open(filepath, 'r', encoding='utf-8', errors='ignore') as f:
                text = f.read()

        elif ext == '.pdf':
            reader = pypdf.PdfReader(filepath)
            for page in reader.pages:
                extracted = page.extract_text()
                if extracted:
                    text += extracted + "\n"
                    for line in extracted.split("\n"):
                        line = line.strip()
                        if line and len(line) < 80 and not line.endswith(("。", "，", "；", ".", ",", ";")):
                            headings.append(line)

        elif ext == '.docx':
            try:
                doc = docx.Document(filepath)
            except KeyError as e:
                print(f"  ⚠️ docx 结构损坏，尝试部分提取 ({os.path.basename(filepath)}): {e}")
                q.put(("", []))
                return
            for para in doc.paragraphs:
                if para.text.strip():
                    text += para.text + "\n"
                    style_name = para.style.name if para.style else ""
                    if style_name.startswith("Heading") or style_name.startswith("标题"):
                        headings.append(para.text.strip())
            try:
                for table in doc.tables:
                    for row in table.rows:
                        cells = set()
                        row_parts = []
                        for cell in row.cells:
                            ct = cell.text.strip()
                            if ct and ct not in cells:
                                cells.add(ct)
                                row_parts.append(ct)
                        if row_parts:
                            text += " | ".join(row_parts) + "\n"
            except KeyError:
                pass

        elif ext == '.xlsx':
            wb = openpyxl.load_workbook(filepath, read_only=True, data_only=True)
            for sheet in wb.worksheets:
                headings.append(sheet.title)
                text += f"Sheet: {sheet.title}\n"
                for row in sheet.iter_rows(values_only=True):
                    row_text = " ".join(str(cell) for cell in row if cell is not None)
                    if row_text.strip():
                        text += row_text + "\n"
            wb.close()

        elif ext == '.pptx':
            prs = Presentation(filepath)
            for slide in prs.slides:
                if slide.shapes.title and slide.shapes.title.text.strip():
                    headings.append(slide.shapes.title.text.strip())
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text += shape.text + "\n"

        elif ext in config.MEDIA_EXTENSIONS:
            pass

    except Exception as e:
        err_type = type(e).__name__
        print(f"  ⚠️ 文件读取异常 ({os.path.basename(filepath)}): [{err_type}] {e}")

    q.put((text.strip(), headings))


def _read_text_direct(filepath: str) -> str:
    """直接读取文本文件（不会挂起，无需子进程）"""
    with open(filepath, 'r', encoding='utf-8', errors='ignore') as f:
        return f.read()


def _read_via_subprocess(filepath: str, ext: str) -> tuple[str, list[str]]:
    """通过子进程读取办公文档，防止 C 扩展死锁。返回 (content, headings)"""
    q = Queue()
    p = Process(target=_read_file_worker, args=(filepath, ext, q))
    p.start()
    p.join(timeout=FILE_READ_TIMEOUT)

    if p.is_alive():
        p.terminate()
        p.join(timeout=5)
        if p.is_alive():
            p.kill()
            p.join()
        filename = normalize_path(os.path.basename(filepath))
        print(f"  ⏰ 超时 ({FILE_READ_TIMEOUT}s): {filename}")
        return "", []

    try:
        return q.get_nowait()
    except Exception:
        return "", []


def _extract_md_headings(content: str) -> list[str]:
    """从 Markdown 内容中提取标题行"""
    return re.findall(r'^#{1,6}\s+(.+)$', content, re.MULTILINE)


def _is_valid_zip(filepath: str) -> bool:
    """Check if file starts with ZIP magic bytes (PK\\x03\\x04)."""
    try:
        with open(filepath, 'rb') as f:
            return f.read(4) == ZIP_MAGIC
    except OSError:
        return False


def load_file_content(filepath: str, ext: str) -> tuple[str, list[str]]:
    """读取文件内容 + 标题。返回 (content, headings)"""
    try:
        size_bytes = os.path.getsize(filepath)
        size_mb = size_bytes / (1024 * 1024)
        if size_mb > MAX_FILE_SIZE_MB:
            print(f"  ⏭️ 跳过 ({size_mb:.0f}MB): {os.path.basename(filepath)}")
            return "", []
    except OSError:
        return "", []

    if ext in config.TEXT_EXTENSIONS:
        try:
            content = _read_text_direct(filepath)
            headings = _extract_md_headings(content) if ext == '.md' else []
            return content, headings
        except Exception:
            return "", []

    if ext in config.MEDIA_EXTENSIONS:
        return "", []

    if ext in config.OFFICE_EXTENSIONS:
        if ext == '.xlsx' and size_mb > XLSX_SKIP_SIZE_MB:
            print(f"  ⏭️ xlsx 过大 ({size_mb:.0f}MB), 仅索引文件名: {os.path.basename(filepath)}")
            return "", []

        if ext in ('.xlsx', '.docx', '.pptx') and not _is_valid_zip(filepath):
            print(f"  ⏭️ 非标准格式(加密/损坏), 仅索引文件名: {os.path.basename(filepath)}")
            return "", []

        return _read_via_subprocess(filepath, ext)

    return "", []


def build_documents_for_file(filepath: str, filename: str, ext: str, source_type: str = "file") -> list[Document]:
    filepath_norm = normalize_path(filepath)
    filename_norm = normalize_path(filename)
    try:
        st = os.stat(filepath)
        mtime = st.st_mtime
        ctime = getattr(st, "st_birthtime", st.st_ctime)
    except OSError:
        mtime = 0
        ctime = 0
    base_meta = {
        "source": filepath_norm,
        "filename": filename_norm,
        "type": ext,
        "source_type": source_type,
        "mtime": mtime,
        "ctime": ctime,
    }

    documents = []

    # 支持多目录：找到文件所属的 target_dir 并计算相对路径（优先匹配最长前缀）
    path_parts = filepath_norm
    for d in sorted(config.get_target_dirs(), key=len, reverse=True):
        if filepath_norm.startswith(d + "/") or filepath_norm == d:
            path_parts = filepath_norm[len(d) :].lstrip("/")
            break
    name_doc_text = _truncate_for_embed(f"文件名: {filename_norm}\n路径: {path_parts}")
    documents.append(Document(page_content=name_doc_text, metadata={**base_meta, "chunk_type": "filename"}))

    content, headings = load_file_content(filepath, ext)

    if headings:
        headings_text = "\n".join(dict.fromkeys(headings))
        heading_content = _truncate_for_embed(f"[{filename_norm}]\n{headings_text}")
        documents.append(Document(
            page_content=heading_content,
            metadata={**base_meta, "chunk_type": "heading"},
        ))

    if not content:
        return documents

    if len(content) > config.MAX_CONTENT_LENGTH:
        content = content[:config.MAX_CONTENT_LENGTH]

    chunks = _splitter.split_text(content)

    for idx, chunk in enumerate(chunks):
        chunk_text = _truncate_for_embed(f"[{filename_norm}]\n{chunk}")
        documents.append(Document(
            page_content=chunk_text,
            metadata={**base_meta, "chunk_type": "content", "chunk_idx": idx},
        ))

    return documents


def scan_files():
    documents = []
    target_dirs = config.get_target_dirs()
    print(f"📂 开始扫描: {target_dirs}")

    if config.INDEX_ONLY_KEYWORDS:
        print(f"🔎 过滤关键词: {config.INDEX_ONLY_KEYWORDS}")

    cache_path = getattr(config, "SCAN_CACHE_PATH", None)
    conn = None
    if cache_path:
        conn = sqlite3.connect(cache_path)
        _init_scan_cache(conn)

    start = time.time()
    file_count = 0
    scanned = 0
    cache_hits = 0

    for target_dir in target_dirs:
        if not os.path.isdir(target_dir):
            print(f"  ⚠️ 跳过不存在的目录: {target_dir}")
            continue
        for root, dirs, files in os.walk(target_dir):
            dirs[:] = [d for d in dirs if not d.startswith('.')]

            for file in files:
                if file.startswith('.'):
                    continue
                scanned += 1

                raw_path = os.path.join(root, file)
                filepath = normalize_path(raw_path)

                if config.INDEX_ONLY_KEYWORDS:
                    if not any(kw in filepath for kw in config.INDEX_ONLY_KEYWORDS):
                        continue

                _, ext = os.path.splitext(file)
                ext = ext.lower()
                if ext not in config.SUPPORTED_EXTENSIONS:
                    continue

                # 大文件提前跳过，避免进入解析流程
                if ext in (config.TEXT_EXTENSIONS | config.OFFICE_EXTENSIONS):
                    try:
                        if os.path.getsize(filepath) > MAX_FILE_SIZE_MB * 1024 * 1024:
                            continue
                    except OSError:
                        continue

                try:
                    mtime = os.path.getmtime(filepath)
                except OSError:
                    mtime = 0

                if conn:
                    cached = _load_cached_docs(conn, filepath, mtime)
                    if cached is not None:
                        documents.extend(cached)
                        file_count += 1
                        cache_hits += 1
                        print(f"  已捕获 {file_count} 个文件 (扫描 {scanned}, 缓存 {cache_hits})...", end='\r')
                        continue

                file_docs = build_documents_for_file(filepath, file, ext)
                documents.extend(file_docs)
                if conn:
                    _save_cached_docs(conn, filepath, mtime, "file", file_docs, auto_commit=False)
                file_count += 1
                if file_count % 50 == 0:
                    if conn:
                        conn.commit()
                    print(f"  已捕获 {file_count} 个文件 (扫描 {scanned})...", end='\r')

    if conn:
        conn.commit()
        conn.close()

    duration = time.time() - start
    hit_info = f", 缓存命中 {cache_hits}" if cache_hits else ""
    print(f"\n✅ 扫描完成: {file_count} 个文件 → {len(documents)} 个文档片段 (扫描 {scanned} 个{hit_info})")
    return documents, duration


def _parse_front_matter(content: str) -> tuple[dict, str]:
    """Parse YAML front matter from markdown content. Returns (metadata, body)."""
    if not content.startswith("---"):
        return {}, content
    end = content.find("\n---", 3)
    if end == -1:
        return {}, content
    try:
        meta = yaml.safe_load(content[3:end]) or {}
    except yaml.YAMLError:
        meta = {}
    body = content[end + 4:].strip()
    return meta, body


def build_documents_for_mweb(filepath: str, content: str) -> list[Document]:
    """Build document chunks for a single MWeb exported markdown note."""
    filepath_norm = normalize_path(filepath)
    meta, body = _parse_front_matter(content)

    title = meta.get("title", "")
    if not title:
        m = re.search(r'^#\s+(.+)$', body, re.MULTILINE)
        title = m.group(1).strip() if m else os.path.splitext(os.path.basename(filepath))[0]

    categories = meta.get("categories", [])
    category_str = categories[0] if categories else ""
    mweb_uuid = str(meta.get("mweb_uuid", ""))
    try:
        st = os.stat(filepath)
        mtime = st.st_mtime
        ctime = getattr(st, "st_birthtime", st.st_ctime)
    except OSError:
        mtime = 0
        ctime = 0

    base_meta = {
        "source": filepath_norm,
        "filename": title,
        "type": ".md",
        "source_type": "mweb",
        "mweb_uuid": mweb_uuid,
        "categories": category_str,
        "mtime": mtime,
        "ctime": ctime,
    }

    documents = []

    name_doc = _truncate_for_embed(f"笔记: {title}\n分类: {category_str}")
    documents.append(Document(page_content=name_doc, metadata={**base_meta, "chunk_type": "filename"}))

    headings = _extract_md_headings(body)
    if headings:
        heading_text = _truncate_for_embed(f"[{title}]\n" + "\n".join(dict.fromkeys(headings)))
        documents.append(Document(page_content=heading_text, metadata={**base_meta, "chunk_type": "heading"}))

    if not body:
        return documents

    if len(body) > config.MAX_CONTENT_LENGTH:
        body = body[:config.MAX_CONTENT_LENGTH]

    chunks = _splitter.split_text(body)
    for idx, chunk in enumerate(chunks):
        chunk_text = _truncate_for_embed(f"[{title}]\n{chunk}")
        documents.append(Document(
            page_content=chunk_text,
            metadata={**base_meta, "chunk_type": "content", "chunk_idx": idx},
        ))

    return documents


def scan_mweb_notes():
    """Scan MWeb exported markdown notes."""
    if not getattr(config, "ENABLE_MWEB", True):
        print("ℹ️ 已关闭 MWeb 数据源（ENABLE_MWEB=False），跳过 MWeb 扫描。")
        return [], 0.0
    mweb_dir = config.MWEB_DIR
    if not os.path.isdir(mweb_dir):
        print(f"⚠️ MWeb 目录不存在: {mweb_dir}")
        return [], 0.0

    cache_path = getattr(config, "SCAN_CACHE_PATH", None)
    conn = None
    if cache_path:
        conn = sqlite3.connect(cache_path)
        _init_scan_cache(conn)

    documents = []
    print(f"📓 开始扫描 MWeb 笔记: {mweb_dir}")
    start = time.time()
    note_count = 0
    cache_hits = 0

    for root, dirs, files in os.walk(mweb_dir):
        dirs[:] = [d for d in dirs if not d.startswith('.')]
        for file in files:
            if not file.endswith('.md') or file.startswith('.'):
                continue
            filepath = normalize_path(os.path.join(root, file))
            try:
                mtime = os.path.getmtime(filepath)
            except OSError:
                mtime = 0

            if conn:
                cached = _load_cached_docs(conn, filepath, mtime)
                if cached is not None:
                    documents.extend(cached)
                    note_count += 1
                    cache_hits += 1
                    if note_count % 50 == 0:
                        print(f"  已捕获 {note_count} 篇笔记 (缓存 {cache_hits})...", end='\r')
                    continue

            try:
                with open(filepath, 'r', encoding='utf-8', errors='ignore') as f:
                    content = f.read()
            except OSError:
                continue

            note_docs = build_documents_for_mweb(filepath, content)
            documents.extend(note_docs)
            if conn:
                _save_cached_docs(conn, filepath, mtime, "mweb", note_docs, auto_commit=False)
            note_count += 1
            if note_count % 50 == 0:
                if conn:
                    conn.commit()
                print(f"  已捕获 {note_count} 篇笔记...", end='\r')

    if conn:
        conn.commit()
        conn.close()

    duration = time.time() - start
    hit_info = f", 缓存命中 {cache_hits}" if cache_hits else ""
    print(f"\n✅ MWeb 扫描完成: {note_count} 篇笔记 → {len(documents)} 个文档片段{hit_info}")
    return documents, duration


def _cleanup_orphaned_hnsw_dirs(client):
    """Remove orphaned HNSW segment directories left by previous index builds."""
    import sqlite3 as _sqlite3
    db_path = os.path.join(config.PERSIST_DIRECTORY, "chroma.sqlite3")
    if not os.path.isfile(db_path):
        return
    try:
        conn = _sqlite3.connect(db_path)
        active_ids = {row[0] for row in conn.execute("SELECT id FROM segments").fetchall()}
        conn.close()
    except Exception:
        return
    for entry in os.listdir(config.PERSIST_DIRECTORY):
        full = os.path.join(config.PERSIST_DIRECTORY, entry)
        if os.path.isdir(full) and entry not in active_ids:
            import shutil
            shutil.rmtree(full, ignore_errors=True)
            print(f"  已清理孤立索引目录: {entry}")


def build_index():
    total_start = time.time()

    docs, scan_time = scan_files()

    mweb_docs, mweb_time = scan_mweb_notes()
    docs.extend(mweb_docs)
    scan_time += mweb_time

    # 清理已删除文件的扫描缓存
    cache_path = getattr(config, "SCAN_CACHE_PATH", None)
    if cache_path and docs:
        valid_sources = {d.metadata.get("source", "") for d in docs if d.metadata.get("source")}
        conn = sqlite3.connect(cache_path)
        _init_scan_cache(conn)
        _prune_scan_cache(conn, valid_sources)
        conn.close()

    if not docs:
        print("❌ 未找到任何文档。请检查 config.py 中的 TARGET_DIR 和 INDEX_ONLY_KEYWORDS。")
        return

    # 最终校验：确保所有 doc 的 page_content 符合 API 限制
    for d in docs:
        if len(d.page_content) > EMBED_MAX_CHARS:
            d.page_content = d.page_content[:EMBED_MAX_CHARS]
        if not d.page_content.strip():
            d.page_content = " "

    print("🗑️ 清除旧索引...")
    import chromadb
    client = chromadb.PersistentClient(path=config.PERSIST_DIRECTORY)
    existing = [c.name for c in client.list_collections()]
    if "local_files" in existing:
        client.delete_collection("local_files")
        print("  已删除旧 collection。")

    _cleanup_orphaned_hnsw_dirs(client)

    print(f"🧠 正在生成向量 (模型: {config.EMBEDDING_MODEL})...")
    embed_start = time.time()

    if not _ensure_dashscope_api_key():
        return
    embeddings = CachedEmbeddings(
        model=config.EMBEDDING_MODEL,
        cache_path=config.EMBEDDING_CACHE_PATH,
    )

    batch_size = 50
    total = len(docs)
    max_retries = 3
    retry_delay = 5

    skip_count = 0
    db = None
    try:
        for i in range(0, total, batch_size):
            batch = docs[i:i + batch_size]
            last_err = None
            for attempt in range(max_retries):
                try:
                    if db is None:
                        db = Chroma.from_documents(
                            documents=batch,
                            embedding=embeddings,
                            persist_directory=config.PERSIST_DIRECTORY,
                            collection_name="local_files",
                            collection_metadata={"hnsw:space": "cosine"},
                        )
                    else:
                        db.add_documents(batch)
                    break
                except Exception as e:
                    last_err = e
                    err_str = str(e)
                    if attempt < max_retries - 1:
                        print(f"\n  ⚠️ 批次 {i}-{i+len(batch)} 失败，{retry_delay}s 后重试 ({attempt+1}/{max_retries}): {err_str[:80]}...")
                        time.sleep(retry_delay)
                    elif "2048" in err_str or "InvalidParameter" in err_str:
                        # 逐条重试，跳过仍失败的文档
                        print(f"\n  ⚠️ 批次失败，逐条重试 (跳过异常文档)...")
                        for j, doc in enumerate(batch):
                            try:
                                if db is None:
                                    db = Chroma.from_documents(
                                        documents=[doc],
                                        embedding=embeddings,
                                        persist_directory=config.PERSIST_DIRECTORY,
                                        collection_name="local_files",
                                        collection_metadata={"hnsw:space": "cosine"},
                                    )
                                else:
                                    db.add_documents([doc])
                            except Exception:
                                skip_count += 1
                                if skip_count <= 5:
                                    print(f"  ⏭️ 跳过文档 (len={len(doc.page_content)}): {doc.page_content[:50]}...")
                        break
                    else:
                        raise last_err

            pct = min((i + batch_size) / total * 100, 100)
            print(f"  进度: {pct:.0f}% ({min(i + batch_size, total)}/{total})", end='\r')

            if i + batch_size < total:
                time.sleep(0.3)

    except Exception as e:
        print(f"\n❌ 向量化错误: {e}")
        return

    embed_time = time.time() - embed_start
    total_time = time.time() - total_start

    print("\n" + "=" * 40)
    print("🎉 索引构建成功！")
    print(f"  文件扫描: {scan_time:.2f}s")
    print(f"  向量化与存储: {embed_time:.2f}s")
    print(f"  总耗时: {total_time:.2f}s")
    print(f"  文档片段: {total}" + (f" (跳过 {skip_count})" if skip_count else ""))
    print(f"  嵌入缓存: {embeddings.stats_str()}")
    print(f"  数据库: {os.path.abspath(config.PERSIST_DIRECTORY)}")
    print("=" * 40)


if __name__ == "__main__":
    build_index()
